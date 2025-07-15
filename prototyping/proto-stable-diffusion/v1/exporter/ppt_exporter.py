from pathlib import Path
import ipdb

from pptx.dml.color import RGBColor
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pathlib import Path
from pptx.dml.color import RGBColor


def build_ppt(templates, slides, title, output_path= "generated_presentation.pptx"):
    prs = Presentation()

    prs = create_title_and_background(prs, templates["title_slide"], slides[0]["title"])

    for slide_data in slides[1:]:
        slide_type = slide_data.get("type", "content")

        # SECTION slide
        if slide_type == "section":
            slide_layout = prs.slide_layouts[0]  # Title Slide
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_data.get("title", "Section")
            continue

        # Default slide layout
        slide_layout = prs.slide_layouts[5]  # Title Only
        slide = prs.slides.add_slide(slide_layout)

        # Add background image for the slide
        add_background_image(slide, templates["slides"], prs.slide_width, prs.slide_height)

        # Title text on top of the background
        slide.shapes.title.text = slide_data.get("title", "Untitled Slide")
       
        has_image = slide_type == "visual" and "image" in slide_data and Path(slide_data["image"]).exists()

        # Text content box — width adjusts based on image presence
        text_width = Inches(4.5) if has_image else Inches(8.5)
        textbox = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), text_width, Inches(4.5))
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        for bullet in slide_data.get("content", []):
            para = text_frame.add_paragraph()
            para.text = f"• {bullet}"
            #para.level = 0
            para.font.size = Pt(16)
            para.font.name = "Arial"
            para.alignment = PP_ALIGN.LEFT 
    
        # Add image on the right (only if exists)
        if has_image:
            slide.shapes.add_picture(
                str(slide_data["image"]),
                left=Inches(5.3),
                top=Inches(1.7),
                width=Inches(4.0)
            )

    prs.save(output_path)
    return Path(output_path)


def create_title_and_background(prs, background_image_path, title_text):
    """
    Create the first slide with a background image and title text.
    Args:
        prs (Presentation): The Presentation object to which the slide will be added.
        background_image_path (str): Path to the background image file.
        title_text (str): Title text to display on the slide.
    Returns:
        Presentation: The updated Presentation object with the title slide added.
    """

    # Add a slide with a blank layout (first slide)
    # 6 is for a blank slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Set the background image for the first slide
    slide.shapes.add_picture(background_image_path, 0, 0,
                             width=prs.slide_width,
                             height=prs.slide_height)

    # Add the title text on top of the background image
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), 
                                         Inches(8), Inches(1))
    title_text_frame = title_box.text_frame
    title_text_frame.text = title_text

    # Customize the title font (optional)
    title_text_frame.paragraphs[0].font.size = Inches(1)
    title_text_frame.paragraphs[0].font.bold = True
    title_text_frame.paragraphs[0].font.color.rgb = RGBColor(16, 16, 16)

    return prs


def add_background_image(slide, image_path, width, height):
    """
    Add a background image to a slide.
    Args:
        slide (pptx.slide.Slide): The slide to which the background image will be added.
        image_path (str): Path to the background image file.
        width (Inches): Width of the slide.
        height (Inches): Height of the slide.
    """
    
    if not Path(image_path).exists():
        raise FileNotFoundError(f"Background image not found at {image_path}")
    
    slide.shapes.add_picture(image_path, 0, 0, width=width, height=height)
