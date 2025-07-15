from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create a presentation
prs = Presentation()

# Add a slide with a blank layout
slide_layout = prs.slide_layouts[5]  # Blank layout
slide = prs.slides.add_slide(slide_layout)

# Add a text box
textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
text_frame = textbox.text_frame
text_frame.word_wrap = True

# Bullet list
content = ["First bullet", "Second bullet", "Third bullet"]

for bullet in content:
    para = text_frame.add_paragraph()

    # Set the text for the paragraph
    para.text = f"• {bullet}"  # Prepend the bullet symbol to the text

    # Customize the font size and alignment
    para.font.size = Pt(16)
    para.font.name = "Arial"
    para.alignment = PP_ALIGN.LEFT

# Save the presentation
prs.save("./data/tmp/test_presentation.pptx")
