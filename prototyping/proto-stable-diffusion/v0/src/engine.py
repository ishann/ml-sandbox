import ipdb
import os
from transformers import pipeline as hf_pipeline
from pptx import Presentation
from pptx.util import Inches
from src.utils import extract_text_from_pdf
from src.utils import generate_image

# Generate presentation from PDF
def create_presentation_from_pdf(pdf_path):
    texts = extract_text_from_pdf(pdf_path)
    prs = Presentation()
    
    # Load summarizer (for slide bullets)
    summarizer = hf_pipeline("summarization", model="facebook/bart-large-cnn")

    ipdb.set_trace()

    for i, text in enumerate(texts):
        short_summary = summarizer(text, max_length=60, min_length=30, do_sample=False)[0]['summary_text']
        image_path = os.path.join(TEMP_IMAGE_DIR, f"slide_{i+1}.png")
        generate_image(short_summary, image_path)

        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only layout
        title_shape = slide.shapes.title
        title_shape.text = short_summary

        left = Inches(1)
        top = Inches(2)
        pic = slide.shapes.add_picture(image_path, left, top, width=Inches(6))

    prs.save(OUTPUT_PPTX)
    print(f"Presentation saved to: {OUTPUT_PPTX}")
